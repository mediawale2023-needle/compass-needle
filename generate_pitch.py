"""
Generate Needle Pitch Deck — needle_pitch_deck.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ─────────────────────────────────────────────
# BRAND COLOURS
# ─────────────────────────────────────────────
NAVY       = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy – slide backgrounds
DARK_CARD  = RGBColor(0x16, 0x27, 0x3A)   # slightly lighter card bg
GOLD       = RGBColor(0xF5, 0xA6, 0x23)   # amber / gold accent
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xCC, 0xD6, 0xE0)
MID_GREY   = RGBColor(0x8A, 0xA0, 0xB4)
GREEN      = RGBColor(0x2E, 0xCC, 0x71)
RED        = RGBColor(0xE7, 0x4C, 0x3C)
TEAL       = RGBColor(0x1A, 0xBC, 0x9C)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # completely blank


# ─────────────────────────────────────────────
# HELPERS
# ─────────────────────────────────────────────

def fill_bg(slide, color: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape

def add_textbox(slide, text, left, top, width, height,
                font_name="Calibri", font_size=14, bold=False, italic=False,
                color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_multiline_textbox(slide, lines, left, top, width, height,
                           font_name="Calibri", font_size=12,
                           color=WHITE, bold=False, line_spacing=1.15):
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    from lxml import etree
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (line_text, line_bold, line_size, line_color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line_text
        run.font.name = font_name
        run.font.size = Pt(line_size or font_size)
        run.font.bold = line_bold if line_bold is not None else bold
        run.font.color.rgb = line_color or color
    return txBox

def accent_bar(slide, top_offset=Inches(0.0), height=Inches(0.06), color=GOLD):
    add_rect(slide, 0, top_offset, SLIDE_W, height, color)

def divider_line(slide, top, color=MID_GREY, height=Inches(0.01)):
    add_rect(slide, Inches(0.6), top, SLIDE_W - Inches(1.2), height, color)

def card(slide, left, top, width, height, fill=DARK_CARD):
    return add_rect(slide, left, top, width, height, fill)

def pill_label(slide, text, left, top, width=Inches(1.8), height=Inches(0.32),
               bg=GOLD, text_color=NAVY):
    add_rect(slide, left, top, width, height, bg)
    add_textbox(slide, text, left, top + Inches(0.04), width, height,
                font_size=9, bold=True, color=text_color, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# SLIDE 1 — HERO / TITLE
# ─────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
fill_bg(s, NAVY)

# Large gold accent bar top
accent_bar(s, Inches(0), Inches(0.08))

# Diagonal gold accent strip bottom-right
add_rect(s, Inches(9.5), Inches(4.5), Inches(4), Inches(3.5), DARK_CARD)
add_rect(s, Inches(10.2), Inches(5.0), Inches(3.0), Inches(2.3), RGBColor(0x1E, 0x35, 0x50))

# Watermark background text
add_textbox(s, "NEEDLE", Inches(5.5), Inches(1.5), Inches(7.5), Inches(5),
            font_size=160, bold=True, color=RGBColor(0x12, 0x24, 0x38),
            align=PP_ALIGN.LEFT)

# Tag / product label
pill_label(s, "PARLIAMENTARY INTELLIGENCE", Inches(0.6), Inches(1.6),
           Inches(3.4), Inches(0.35), bg=GOLD, text_color=NAVY)

# Title
add_textbox(s, "Needle", Inches(0.6), Inches(2.1), Inches(7), Inches(1.6),
            font_size=72, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# Tagline
add_textbox(s, "AI-powered constituency management\nfor every Member of Parliament.",
            Inches(0.6), Inches(3.65), Inches(6.8), Inches(1.2),
            font_size=20, bold=False, color=LIGHT_GREY, align=PP_ALIGN.LEFT)

# Sub-tagline
add_textbox(s, "WhatsApp  ·  AI  ·  Dashboard  ·  Drafting  ·  Intelligence",
            Inches(0.6), Inches(4.8), Inches(8), Inches(0.5),
            font_size=12, bold=False, color=MID_GREY, align=PP_ALIGN.LEFT)

# Bottom accent bar
accent_bar(s, SLIDE_H - Inches(0.08), Inches(0.08))

# Bottom right branding
add_textbox(s, "MediaWale 2023  ·  Proprietary & Confidential",
            Inches(0), SLIDE_H - Inches(0.55), SLIDE_W - Inches(0.3), Inches(0.4),
            font_size=9, color=MID_GREY, align=PP_ALIGN.RIGHT)


# ─────────────────────────────────────────────
# SLIDE 2 — THE PROBLEM
# ─────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
fill_bg(s, NAVY)
accent_bar(s, Inches(0), Inches(0.08))

add_textbox(s, "THE PROBLEM", Inches(0.6), Inches(0.25), Inches(4), Inches(0.4),
            font_size=10, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

add_textbox(s, "Indian MPs are drowning\nin their own constituency.",
            Inches(0.6), Inches(0.65), Inches(8), Inches(1.4),
            font_size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

divider_line(s, Inches(2.1))

# Pain-point cards — 3 across
pain_points = [
    ("1,000+\nGrievances / Month",
     "Most MPs receive thousands of WhatsApp messages, letters, and phone calls — with zero system to track or resolve them."),
    ("0\nDigital Tools Built for MPs",
     "There is no Indian-made, parliament-aware SaaS product. MPs rely on personal assistants, Excel sheets, and WhatsApp groups."),
    ("₹1,000 Cr+\nCSR Funds Untapped",
     "MPs struggle to connect their constituents with available corporate CSR money. The matching is manual, slow, and opaque."),
]

col_w = Inches(3.8)
col_gap = Inches(0.3)
start_x = Inches(0.5)

for i, (title, desc) in enumerate(pain_points):
    cx = start_x + i * (col_w + col_gap)
    cy = Inches(2.3)
    card(s, cx, cy, col_w, Inches(3.8))
    add_rect(s, cx, cy, Inches(0.05), Inches(3.8), RED)  # red left edge
    add_textbox(s, title, cx + Inches(0.2), cy + Inches(0.25),
                col_w - Inches(0.4), Inches(1.1),
                font_size=22, bold=True, color=WHITE)
    add_textbox(s, desc, cx + Inches(0.2), cy + Inches(1.4),
                col_w - Inches(0.4), Inches(2.1),
                font_size=11, color=LIGHT_GREY)

# Bottom quote
add_textbox(s, '"Indian democracy\'s last mile is broken. Needle fixes it."',
            Inches(0.6), Inches(6.5), Inches(9), Inches(0.6),
            font_size=13, italic=True, color=GOLD, align=PP_ALIGN.LEFT)

accent_bar(s, SLIDE_H - Inches(0.08), Inches(0.08))


# ─────────────────────────────────────────────
# SLIDE 3 — THE SOLUTION
# ─────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
fill_bg(s, NAVY)
accent_bar(s, Inches(0), Inches(0.08))

add_textbox(s, "THE SOLUTION", Inches(0.6), Inches(0.25), Inches(4), Inches(0.4),
            font_size=10, bold=True, color=GOLD)

add_textbox(s, "Needle — All-in-one AI command\ncenter for the modern MP.",
            Inches(0.6), Inches(0.65), Inches(8), Inches(1.4),
            font_size=36, bold=True, color=WHITE)

divider_line(s, Inches(2.1))

# Central flow diagram — boxes with arrows
flow = [
    ("Citizens", "WhatsApp\nmessage"),
    ("Needle AI", "Classifies\n& routes"),
    ("MP Dashboard", "Manages\ncases"),
    ("Action", "Resolves\n& responds"),
]

box_w = Inches(2.4)
box_h = Inches(1.5)
gap   = Inches(0.5)
total = len(flow) * box_w + (len(flow) - 1) * gap
start = (SLIDE_W - total) / 2
ty = Inches(2.5)

for i, (title, sub) in enumerate(flow):
    bx = start + i * (box_w + gap)
    is_needle = (i == 1)
    bg_col = GOLD if is_needle else DARK_CARD
    tc_col = NAVY if is_needle else WHITE
    sc_col = NAVY if is_needle else MID_GREY
    card(s, bx, ty, box_w, box_h, bg_col)
    add_textbox(s, title, bx, ty + Inches(0.25), box_w, Inches(0.5),
                font_size=15, bold=True, color=tc_col, align=PP_ALIGN.CENTER)
    add_textbox(s, sub, bx, ty + Inches(0.75), box_w, Inches(0.6),
                font_size=10, color=sc_col, align=PP_ALIGN.CENTER)
    if i < len(flow) - 1:
        ax = bx + box_w + Inches(0.05)
        add_textbox(s, "▶", ax, ty + Inches(0.55), Inches(0.4), Inches(0.4),
                    font_size=18, color=GOLD, align=PP_ALIGN.CENTER)

# Feature pills below
features = [
    "WhatsApp Grievance Intake",
    "GPT-4 AI Classification",
    "Multi-language Support",
    "AI Letter & Speech Drafter",
    "Document Copilot",
    "CSR Project Discovery",
    "Constituency News Feed",
    "Analytics & Reporting",
]

pill_w = Inches(2.6)
pill_h = Inches(0.38)
pills_per_row = 4
rx = Inches(0.55)
ry = Inches(4.35)
for i, feat in enumerate(features):
    col_i = i % pills_per_row
    row_i = i // pills_per_row
    px = rx + col_i * (pill_w + Inches(0.18))
    py = ry + row_i * (pill_h + Inches(0.15))
    add_rect(s, px, py, pill_w, pill_h, DARK_CARD)
    add_rect(s, px, py, Inches(0.04), pill_h, GOLD)
    add_textbox(s, feat, px + Inches(0.12), py + Inches(0.06),
                pill_w - Inches(0.15), pill_h - Inches(0.1),
                font_size=10, bold=False, color=LIGHT_GREY)

accent_bar(s, SLIDE_H - Inches(0.08), Inches(0.08))


# ─────────────────────────────────────────────
# SLIDE 4 — WHATSAPP GRIEVANCE ENGINE
# ─────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
fill_bg(s, NAVY)
accent_bar(s, Inches(0), Inches(0.08))

add_textbox(s, "FEATURE 01 — GRIEVANCE INTAKE", Inches(0.6), Inches(0.25), Inches(6), Inches(0.4),
            font_size=10, bold=True, color=GOLD)
add_textbox(s, "Every citizen on WhatsApp.\nEvery complaint in one dashboard.",
            Inches(0.6), Inches(0.65), Inches(8), Inches(1.4),
            font_size=34, bold=True, color=WHITE)

divider_line(s, Inches(2.1))

# Left column — how it works
steps = [
    ("01", "Citizen sends a message in any language — Hindi, Marathi, Kannada, English or Hinglish."),
    ("02", "Needle's AI (GPT-4) reads it, detects the language, classifies the grievance into 20 categories."),
    ("03", "Geography engine maps the citizen to the correct Assembly Constituency automatically."),
    ("04", "Case is logged with all metadata. AI sends an instant, personalised acknowledgment."),
    ("05", "MP sees a clean dashboard: status, category, location, priority — ready to act."),
]
sy = Inches(2.3)
for num, text in steps:
    add_rect(s, Inches(0.6), sy, Inches(0.45), Inches(0.45), GOLD)
    add_textbox(s, num, Inches(0.6), sy + Inches(0.04), Inches(0.45), Inches(0.4),
                font_size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_textbox(s, text, Inches(1.2), sy, Inches(5.5), Inches(0.5),
                font_size=10, color=LIGHT_GREY)
    sy += Inches(0.7)

# Right column — stats card
stat_bg = card(s, Inches(7.5), Inches(2.2), Inches(5.3), Inches(4.5))
add_textbox(s, "What this means for MPs",
            Inches(7.7), Inches(2.4), Inches(4.9), Inches(0.4),
            font_size=13, bold=True, color=GOLD)

stats = [
    ("20", "Grievance categories tracked"),
    ("5", "Languages supported"),
    ("543", "Lok Sabha constituencies covered"),
    ("< 3s", "AI response time to citizen"),
    ("100%", "Automatic geography tagging"),
]
sy2 = Inches(2.95)
for val, lbl in stats:
    add_textbox(s, val, Inches(7.7), sy2, Inches(1.3), Inches(0.45),
                font_size=20, bold=True, color=WHITE)
    add_textbox(s, lbl, Inches(9.1), sy2 + Inches(0.08), Inches(3.3), Inches(0.35),
                font_size=10, color=MID_GREY)
    sy2 += Inches(0.65)

accent_bar(s, SLIDE_H - Inches(0.08), Inches(0.08))


# ─────────────────────────────────────────────
# SLIDE 5 — AI DRAFTER
# ─────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
fill_bg(s, NAVY)
accent_bar(s, Inches(0), Inches(0.08))

add_textbox(s, "FEATURE 02 — AI DRAFTER", Inches(0.6), Inches(0.25), Inches(6), Inches(0.4),
            font_size=10, bold=True, color=GOLD)
add_textbox(s, "From idea to Parliament-ready\ndocument in under 60 seconds.",
            Inches(0.6), Inches(0.65), Inches(9), Inches(1.4),
            font_size=34, bold=True, color=WHITE)

divider_line(s, Inches(2.1))

doc_types = [
    ("Official Letters", "Government of India standard format, F.No., formal Rajbhasha Hindi or English. Ministers, departments, officials."),
    ("Parliament Questions", "Starred & unstarred questions in Lok Sabha / Rajya Sabha format with proper subject lines."),
    ("Private Member Bills", "Structured PMB drafts with clause-by-clause format, statement of objects, and reasons."),
    ("Speeches", "Constituency speeches, budget responses, committee statements with data backing."),
]

dw = Inches(2.9)
dh = Inches(3.5)
dx0 = Inches(0.5)
dy0 = Inches(2.3)

for i, (dtype, desc) in enumerate(doc_types):
    dx = dx0 + i * (dw + Inches(0.25))
    card(s, dx, dy0, dw, dh)
    # top accent
    add_rect(s, dx, dy0, dw, Inches(0.06), GOLD)
    add_textbox(s, dtype, dx + Inches(0.15), dy0 + Inches(0.18),
                dw - Inches(0.3), Inches(0.55),
                font_size=14, bold=True, color=WHITE)
    add_textbox(s, desc, dx + Inches(0.15), dy0 + Inches(0.85),
                dw - Inches(0.3), Inches(2.4),
                font_size=10, color=LIGHT_GREY)

# Bottom row — capabilities
caps = [
    "Tone selector: Assertive · Diplomatic · Formal · Urgent",
    "DNA style templates: Learn the MP's writing voice",
    "Formal Rajbhasha Hindi with correct honorifics",
    "1-click save to Archives",
]
cy2 = Inches(5.95)
for cap in caps:
    add_textbox(s, "✓  " + cap, Inches(0.7), cy2, Inches(12), Inches(0.35),
                font_size=10, color=GREEN)
    cy2 += Inches(0.32)

accent_bar(s, SLIDE_H - Inches(0.08), Inches(0.08))


# ─────────────────────────────────────────────
# SLIDE 6 — CONSTITUENCY INTELLIGENCE
# ─────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
fill_bg(s, NAVY)
accent_bar(s, Inches(0), Inches(0.08))

add_textbox(s, "FEATURE 03 — CONSTITUENCY INTELLIGENCE", Inches(0.6), Inches(0.25), Inches(8), Inches(0.4),
            font_size=10, bold=True, color=GOLD)
add_textbox(s, "Know your constituency better\nthan any opponent ever could.",
            Inches(0.6), Inches(0.65), Inches(8.5), Inches(1.4),
            font_size=34, bold=True, color=WHITE)

divider_line(s, Inches(2.1))

intel_modules = [
    ("Document Copilot",
     TEAL,
     ["Upload any PDF — government report, budget, scheme notification.",
      "Ask questions in plain English or Hindi.",
      "Get instant, cited summaries, risk flags, and action recommendations."]),
    ("CSR Project Discovery",
     GOLD,
     ["Browse 10+ shovel-ready projects: schools, borewells, health centres, roads.",
      "Generate a ready-to-send CSR proposal PDF automatically.",
      "Match constituency needs with corporate Schedule VII funding."]),
    ("News Intelligence",
     RGBColor(0x95, 0x44, 0xF4),
     ["Real-time Google News feed filtered for the MP's constituency.",
      "Constituency-specific + national political news in one view.",
      "Stay ahead of emerging issues before they become grievances."]),
]

mw = Inches(3.8)
mh = Inches(4.0)
mx0 = Inches(0.5)
my0 = Inches(2.25)

for i, (title, accent, bullets) in enumerate(intel_modules):
    mx = mx0 + i * (mw + Inches(0.27))
    card(s, mx, my0, mw, mh)
    add_rect(s, mx, my0, Inches(0.06), mh, accent)
    add_textbox(s, title, mx + Inches(0.2), my0 + Inches(0.2),
                mw - Inches(0.3), Inches(0.5),
                font_size=15, bold=True, color=WHITE)
    by = my0 + Inches(0.9)
    for bullet in bullets:
        add_textbox(s, "▸  " + bullet, mx + Inches(0.2), by,
                    mw - Inches(0.3), Inches(0.7),
                    font_size=10, color=LIGHT_GREY)
        by += Inches(0.85)

accent_bar(s, SLIDE_H - Inches(0.08), Inches(0.08))


# ─────────────────────────────────────────────
# SLIDE 7 — FOR THE MP: BEFORE vs AFTER
# ─────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
fill_bg(s, NAVY)
accent_bar(s, Inches(0), Inches(0.08))

add_textbox(s, "THE MP'S EXPERIENCE", Inches(0.6), Inches(0.25), Inches(6), Inches(0.4),
            font_size=10, bold=True, color=GOLD)
add_textbox(s, "Before Needle vs. After Needle",
            Inches(0.6), Inches(0.65), Inches(10), Inches(0.9),
            font_size=36, bold=True, color=WHITE)

divider_line(s, Inches(1.65))

# Before column
add_rect(s, Inches(0.5), Inches(1.8), Inches(5.8), Inches(0.55), RED)
add_textbox(s, "BEFORE — Without Needle",
            Inches(0.5), Inches(1.8), Inches(5.8), Inches(0.55),
            font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

before_items = [
    "Thousands of WhatsApp messages with no system",
    "PA manually notes complaints in registers / Excel",
    "Hours wasted drafting letters from scratch",
    "No idea which grievances are pending",
    "Missed CSR funding opportunities",
    "No data-driven view of constituency health",
    "Dependent on party networks for news",
    "Reactive to problems rather than proactive",
]
by = Inches(2.5)
for item in before_items:
    card(s, Inches(0.5), by, Inches(5.8), Inches(0.42), RGBColor(0x1F, 0x10, 0x10))
    add_textbox(s, "✗  " + item, Inches(0.65), by + Inches(0.07),
                Inches(5.5), Inches(0.3), font_size=10, color=RGBColor(0xFF, 0xAA, 0xAA))
    by += Inches(0.5)

# After column
add_rect(s, Inches(7.0), Inches(1.8), Inches(5.8), Inches(0.55), GREEN)
add_textbox(s, "AFTER — With Needle",
            Inches(7.0), Inches(1.8), Inches(5.8), Inches(0.55),
            font_size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

after_items = [
    "All citizen messages auto-captured & categorised",
    "AI handles intake 24/7, PA focuses on resolution",
    "Professional letters drafted in under 60 seconds",
    "Real-time case tracker with status & priority",
    "CSR discovery + proposal generation built in",
    "Analytics dashboard: category trends, resolution rate",
    "Constituency-filtered live news in one click",
    "AI flags emerging issues before they escalate",
]
ay = Inches(2.5)
for item in after_items:
    card(s, Inches(7.0), ay, Inches(5.8), Inches(0.42), RGBColor(0x0D, 0x1F, 0x10))
    add_textbox(s, "✓  " + item, Inches(7.15), ay + Inches(0.07),
                Inches(5.5), Inches(0.3), font_size=10, color=RGBColor(0xAA, 0xFF, 0xCC))
    ay += Inches(0.5)

accent_bar(s, SLIDE_H - Inches(0.08), Inches(0.08))


# ─────────────────────────────────────────────
# SLIDE 8 — MARKET OPPORTUNITY
# ─────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
fill_bg(s, NAVY)
accent_bar(s, Inches(0), Inches(0.08))

add_textbox(s, "MARKET OPPORTUNITY", Inches(0.6), Inches(0.25), Inches(6), Inches(0.4),
            font_size=10, bold=True, color=GOLD)
add_textbox(s, "The largest untapped GovTech\nmarket in the world.",
            Inches(0.6), Inches(0.65), Inches(8), Inches(1.4),
            font_size=36, bold=True, color=WHITE)

divider_line(s, Inches(2.1))

# Big numbers
big_stats = [
    ("543", "Lok Sabha MPs", "India's Lower House"),
    ("245", "Rajya Sabha MPs", "India's Upper House"),
    ("4,100+", "State MLAs", "Vidhansabha members"),
    ("1M+", "Local Ward Reps", "Zila Panchayat & ULBs"),
]

bw = Inches(2.9)
bh = Inches(2.1)
bx0 = Inches(0.5)
by0 = Inches(2.3)

for i, (val, title, sub) in enumerate(big_stats):
    bx = bx0 + i * (bw + Inches(0.25))
    card(s, bx, by0, bw, bh)
    add_rect(s, bx, by0, bw, Inches(0.06), GOLD)
    add_textbox(s, val, bx, by0 + Inches(0.2), bw, Inches(0.75),
                font_size=38, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_textbox(s, title, bx, by0 + Inches(1.0), bw, Inches(0.4),
                font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, sub, bx, by0 + Inches(1.45), bw, Inches(0.4),
                font_size=9, color=MID_GREY, align=PP_ALIGN.CENTER)

# TAM SAM SOM
add_textbox(s, "Market Sizing", Inches(0.6), Inches(4.6), Inches(4), Inches(0.4),
            font_size=12, bold=True, color=GOLD)

market_tiers = [
    ("TAM", "₹1,200 Cr / yr", "All elected reps in India @ ₹10K/mo"),
    ("SAM", "₹120 Cr / yr", "All 543 + 245 national MPs @ ₹10K/mo"),
    ("SOM (Y3)", "₹18 Cr / yr", "150 MPs @ ₹10K/mo — 20% Lok Sabha penetration"),
]

for i, (tier, size, desc) in enumerate(market_tiers):
    tx = Inches(0.6) + i * Inches(4.0)
    card(s, tx, Inches(5.1), Inches(3.7), Inches(1.7))
    add_textbox(s, tier, tx + Inches(0.2), Inches(5.2), Inches(0.6), Inches(0.4),
                font_size=11, bold=True, color=GOLD)
    add_textbox(s, size, tx + Inches(0.2), Inches(5.55), Inches(3.3), Inches(0.45),
                font_size=18, bold=True, color=WHITE)
    add_textbox(s, desc, tx + Inches(0.2), Inches(5.95), Inches(3.3), Inches(0.6),
                font_size=9, color=MID_GREY)

accent_bar(s, SLIDE_H - Inches(0.08), Inches(0.08))


# ─────────────────────────────────────────────
# SLIDE 9 — PRODUCT ARCHITECTURE (HOW IT WORKS)
# ─────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
fill_bg(s, NAVY)
accent_bar(s, Inches(0), Inches(0.08))

add_textbox(s, "TECHNICAL ARCHITECTURE", Inches(0.6), Inches(0.25), Inches(6), Inches(0.4),
            font_size=10, bold=True, color=GOLD)
add_textbox(s, "Enterprise-grade. Cloud-native.\nBuilt for scale.",
            Inches(0.6), Inches(0.65), Inches(8), Inches(1.4),
            font_size=34, bold=True, color=WHITE)

divider_line(s, Inches(2.1))

# Three-tier architecture diagram
tiers = [
    ("INPUT LAYER", TEAL, ["WhatsApp (Twilio)", "Citizen Messages", "Multi-language NLP"]),
    ("AI ENGINE", GOLD, ["GPT-4 Classification", "Geography Resolver", "Rajbhasha Drafter", "Document Copilot"]),
    ("DELIVERY LAYER", RGBColor(0x95, 0x44, 0xF4), ["MP Dashboard (Next.js)", "Admin Dashboard", "Case Analytics", "CSR & News Modules"]),
]

tw = Inches(3.6)
th = Inches(3.8)
tx0 = Inches(0.6)
ty0 = Inches(2.3)

for i, (tier_name, accent_color, items) in enumerate(tiers):
    tx = tx0 + i * (tw + Inches(0.45))
    card(s, tx, ty0, tw, th)
    add_rect(s, tx, ty0, tw, Inches(0.06), accent_color)
    add_textbox(s, tier_name, tx, ty0 + Inches(0.15), tw, Inches(0.4),
                font_size=10, bold=True, color=accent_color, align=PP_ALIGN.CENTER)
    iy = ty0 + Inches(0.65)
    for item in items:
        add_rect(s, tx + Inches(0.2), iy, tw - Inches(0.4), Inches(0.45),
                 RGBColor(0x1A, 0x30, 0x45))
        add_textbox(s, item, tx + Inches(0.3), iy + Inches(0.07),
                    tw - Inches(0.5), Inches(0.32),
                    font_size=10, color=WHITE)
        iy += Inches(0.6)

    if i < len(tiers) - 1:
        ax = tx + tw + Inches(0.1)
        add_textbox(s, "→", ax, ty0 + Inches(1.6), Inches(0.35), Inches(0.4),
                    font_size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# Tech stack bar
add_textbox(s, "TECH STACK:", Inches(0.6), Inches(6.3), Inches(1.5), Inches(0.4),
            font_size=9, bold=True, color=MID_GREY)
tech = ["FastAPI  Python 3.11", "Next.js 15", "PostgreSQL  Railway", "OpenAI GPT-4", "Twilio WhatsApp", "Docker  Railway Cloud"]
tx2 = Inches(2.2)
for t in tech:
    add_rect(s, tx2, Inches(6.3), Inches(1.75), Inches(0.35), DARK_CARD)
    add_textbox(s, t, tx2 + Inches(0.1), Inches(6.33), Inches(1.6), Inches(0.3),
                font_size=8, bold=True, color=LIGHT_GREY)
    tx2 += Inches(1.85)

accent_bar(s, SLIDE_H - Inches(0.08), Inches(0.08))


# ─────────────────────────────────────────────
# SLIDE 10 — COMPETITIVE ADVANTAGE / MOAT
# ─────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
fill_bg(s, NAVY)
accent_bar(s, Inches(0), Inches(0.08))

add_textbox(s, "COMPETITIVE ADVANTAGE", Inches(0.6), Inches(0.25), Inches(6), Inches(0.4),
            font_size=10, bold=True, color=GOLD)
add_textbox(s, "Deep moats. First-mover.\nImpossible to replicate quickly.",
            Inches(0.6), Inches(0.65), Inches(9), Inches(1.4),
            font_size=34, bold=True, color=WHITE)

divider_line(s, Inches(2.1))

moats = [
    ("Parliamentary\nDomain Expertise",
     GOLD,
     "Built-in knowledge of Parliament rules, GoI letter formats, Rajbhasha standards, and grievance taxonomy — not generic AI."),
    ("Geography\nData Layer",
     TEAL,
     "5,000+ polling stations mapped to Assembly Constituencies. This data layer took months to build and is proprietary."),
    ("Multi-language\nPolitical NLP",
     RGBColor(0x95, 0x44, 0xF4),
     "AI tuned for Indian political language — Hinglish, regional dialects, government vocabulary across 5+ languages."),
    ("Network\nEffects",
     GREEN,
     "Each onboarded MP generates data that makes the AI smarter for all MPs. Grievance taxonomy improves with every case."),
]

mw2 = Inches(2.9)
mh2 = Inches(3.5)
mx2 = Inches(0.5)
my2 = Inches(2.3)

for i, (title, accent, desc) in enumerate(moats):
    mx = mx2 + i * (mw2 + Inches(0.3))
    card(s, mx, my2, mw2, mh2)
    add_rect(s, mx, my2, mw2, Inches(0.06), accent)
    add_textbox(s, title, mx, my2 + Inches(0.2), mw2, Inches(0.8),
                font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, desc, mx + Inches(0.15), my2 + Inches(1.1),
                mw2 - Inches(0.3), Inches(2.1),
                font_size=10, color=LIGHT_GREY)

# Bottom row
add_textbox(s,
    "No comparable product exists in India for parliamentarians. Needle is building a category.",
    Inches(0.6), Inches(6.05), Inches(12), Inches(0.55),
    font_size=13, italic=True, bold=False, color=GOLD, align=PP_ALIGN.CENTER)

accent_bar(s, SLIDE_H - Inches(0.08), Inches(0.08))


# ─────────────────────────────────────────────
# SLIDE 11 — BUSINESS MODEL
# ─────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
fill_bg(s, NAVY)
accent_bar(s, Inches(0), Inches(0.08))

add_textbox(s, "BUSINESS MODEL", Inches(0.6), Inches(0.25), Inches(6), Inches(0.4),
            font_size=10, bold=True, color=GOLD)
add_textbox(s, "Subscription SaaS with high retention\nand predictable recurring revenue.",
            Inches(0.6), Inches(0.65), Inches(9), Inches(1.4),
            font_size=34, bold=True, color=WHITE)

divider_line(s, Inches(2.1))

# Pricing tiers
tiers_bm = [
    ("Starter", "₹5,000 / mo", ["WhatsApp Intake (500 cases/mo)", "Case Dashboard", "Basic Drafter", "Email Support"]),
    ("Pro", "₹10,000 / mo", ["Unlimited Cases", "Full AI Drafter Suite", "Document Copilot", "CSR Discovery", "Constituency News", "Priority Support"]),
    ("Elite", "₹25,000 / mo", ["Everything in Pro", "Custom DNA Templates", "Multi-seat access", "Dedicated Onboarding", "SLA 99.9%", "WhatsApp Support"]),
]

tw2 = Inches(3.6)
th2 = Inches(4.2)
tx2b = Inches(0.65)
ty2 = Inches(2.2)

for i, (tier_name, price, features) in enumerate(tiers_bm):
    tx = tx2b + i * (tw2 + Inches(0.35))
    is_pro = (i == 1)
    bg_col = DARK_CARD
    card(s, tx, ty2, tw2, th2, bg_col)
    top_bar_col = GOLD if is_pro else MID_GREY
    add_rect(s, tx, ty2, tw2, Inches(0.07), top_bar_col)
    if is_pro:
        pill_label(s, "MOST POPULAR", tx + Inches(0.85), ty2 + Inches(0.15),
                   Inches(1.9), Inches(0.28), bg=GOLD, text_color=NAVY)
    add_textbox(s, tier_name, tx, ty2 + Inches(0.55), tw2, Inches(0.45),
                font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, price, tx, ty2 + Inches(1.0), tw2, Inches(0.55),
                font_size=24, bold=True, color=GOLD if is_pro else LIGHT_GREY,
                align=PP_ALIGN.CENTER)
    divider_line(s, ty2 + Inches(1.6), color=DARK_CARD)
    fy = ty2 + Inches(1.7)
    for feat in features:
        add_textbox(s, "✓  " + feat, tx + Inches(0.2), fy, tw2 - Inches(0.3), Inches(0.35),
                    font_size=9, color=LIGHT_GREY if not is_pro else WHITE)
        fy += Inches(0.38)

# Revenue projection note
add_textbox(s, "At 150 MPs on Pro plan:  150 × ₹10,000 = ₹1.5 Cr / month  =  ₹18 Cr ARR",
            Inches(0.6), Inches(6.55), Inches(12), Inches(0.55),
            font_size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

accent_bar(s, SLIDE_H - Inches(0.08), Inches(0.08))


# ─────────────────────────────────────────────
# SLIDE 12 — TRACTION & ROADMAP
# ─────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
fill_bg(s, NAVY)
accent_bar(s, Inches(0), Inches(0.08))

add_textbox(s, "TRACTION & ROADMAP", Inches(0.6), Inches(0.25), Inches(6), Inches(0.4),
            font_size=10, bold=True, color=GOLD)
add_textbox(s, "Built. Deployed. Ready to scale.",
            Inches(0.6), Inches(0.65), Inches(9), Inches(0.9),
            font_size=36, bold=True, color=WHITE)

divider_line(s, Inches(1.6))

# Left: what's done
add_textbox(s, "PRODUCT MILESTONE", Inches(0.6), Inches(1.8), Inches(5.5), Inches(0.4),
            font_size=10, bold=True, color=GOLD)

milestones = [
    ("LIVE", GREEN, "WhatsApp grievance intake — production ready"),
    ("LIVE", GREEN, "GPT-4 AI classification engine — multi-language"),
    ("LIVE", GREEN, "Geography engine — 543 constituencies mapped"),
    ("LIVE", GREEN, "MP Dashboard — Next.js, JWT auth, full RBAC"),
    ("LIVE", GREEN, "AI Drafter — letters, speeches, PMBs, questions"),
    ("LIVE", GREEN, "Admin command centre — multi-tenant management"),
    ("LIVE", GREEN, "Document Copilot — PDF Q&A for MPs"),
    ("LIVE", GREEN, "CSR discovery & proposal generation"),
    ("LIVE", GREEN, "Constituency News Intelligence feed"),
]

my3 = Inches(2.3)
for status, status_col, text in milestones:
    add_rect(s, Inches(0.6), my3 + Inches(0.06), Inches(0.5), Inches(0.28), status_col)
    add_textbox(s, status, Inches(0.6), my3 + Inches(0.06), Inches(0.5), Inches(0.28),
                font_size=7, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_textbox(s, text, Inches(1.22), my3, Inches(5.0), Inches(0.38),
                font_size=10, color=LIGHT_GREY)
    my3 += Inches(0.46)

# Right: roadmap
add_textbox(s, "GROWTH ROADMAP", Inches(7.2), Inches(1.8), Inches(5.5), Inches(0.4),
            font_size=10, bold=True, color=GOLD)

roadmap = [
    ("Q2 2025", "Onboard first 10 paying MPs — pilot programme"),
    ("Q3 2025", "State assembly (MLA) version — expand to 4,100 seat market"),
    ("Q4 2025", "Mobile app for MPs — iOS & Android"),
    ("Q1 2026", "Vernacular voice grievance intake (IVR + WhatsApp voice)"),
    ("Q2 2026", "State government API integrations — live scheme status"),
    ("Q3 2026", "100 MPs milestone — Series A fundraise"),
    ("2027", "Expand to Sri Lanka, Bangladesh, Nepal — SAARC play"),
]

ry2 = Inches(2.3)
for quarter, desc in roadmap:
    add_rect(s, Inches(7.2), ry2 + Inches(0.05), Inches(1.1), Inches(0.3), DARK_CARD)
    add_textbox(s, quarter, Inches(7.2), ry2 + Inches(0.05), Inches(1.1), Inches(0.3),
                font_size=8, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_textbox(s, desc, Inches(8.4), ry2, Inches(4.4), Inches(0.4),
                font_size=10, color=LIGHT_GREY)
    ry2 += Inches(0.52)

accent_bar(s, SLIDE_H - Inches(0.08), Inches(0.08))


# ─────────────────────────────────────────────
# SLIDE 13 — THE ASK (CTA)
# ─────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
fill_bg(s, NAVY)
accent_bar(s, Inches(0), Inches(0.08))

# Background decoration
add_rect(s, Inches(8.5), Inches(1.5), Inches(5), Inches(6), DARK_CARD)
add_rect(s, Inches(9.2), Inches(2.2), Inches(4.0), Inches(4.5), RGBColor(0x1A, 0x2E, 0x42))

add_textbox(s, "THE ASK", Inches(0.6), Inches(0.25), Inches(4), Inches(0.4),
            font_size=10, bold=True, color=GOLD)

add_textbox(s, "Join the future\nof Indian democracy.",
            Inches(0.6), Inches(0.75), Inches(7.5), Inches(1.6),
            font_size=40, bold=True, color=WHITE)

divider_line(s, Inches(2.45))

# For MPs
add_rect(s, Inches(0.6), Inches(2.6), Inches(5.8), Inches(0.5), RGBColor(0x14, 0x2A, 0x1E))
add_rect(s, Inches(0.6), Inches(2.6), Inches(0.06), Inches(0.5), GREEN)
add_textbox(s, "FOR MEMBERS OF PARLIAMENT", Inches(0.72), Inches(2.67), Inches(5.6), Inches(0.35),
            font_size=11, bold=True, color=GREEN)

mp_asks = [
    "Start your free 30-day pilot — no IT setup required",
    "Get your constituency live on Needle in under 48 hours",
    "Join a cohort of forward-thinking MPs who put constituents first",
]
ay2 = Inches(3.25)
for ask in mp_asks:
    add_textbox(s, "▸  " + ask, Inches(0.7), ay2, Inches(5.6), Inches(0.38),
                font_size=11, color=LIGHT_GREY)
    ay2 += Inches(0.42)

# For Investors
iy2 = Inches(4.8)
add_rect(s, Inches(0.6), iy2, Inches(5.8), Inches(0.5), RGBColor(0x14, 0x22, 0x2A))
add_rect(s, Inches(0.6), iy2, Inches(0.06), Inches(0.5), GOLD)
add_textbox(s, "FOR INVESTORS", Inches(0.72), iy2 + Inches(0.07), Inches(5.6), Inches(0.35),
            font_size=11, bold=True, color=GOLD)

investor_asks = [
    "Raising ₹2 Cr seed to fuel MP acquisition & team growth",
    "Use of funds: 50% sales, 30% product, 20% ops",
    "Targeting 50 MPs in Y1 → ₹6 Cr ARR → Series A in 24 months",
]
iy2b = Inches(5.45)
for ask in investor_asks:
    add_textbox(s, "▸  " + ask, Inches(0.7), iy2b, Inches(5.6), Inches(0.38),
                font_size=11, color=LIGHT_GREY)
    iy2b += Inches(0.42)

# Right panel — contact / closing
add_textbox(s, "needle", Inches(9.0), Inches(2.6), Inches(3.5), Inches(0.8),
            font_size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, "Parliamentary Intelligence Platform",
            Inches(8.8), Inches(3.4), Inches(3.8), Inches(0.45),
            font_size=10, color=MID_GREY, align=PP_ALIGN.CENTER)

divider_line(s, Inches(4.0))

add_textbox(s, "contact@mediawale.in",
            Inches(8.8), Inches(4.15), Inches(3.8), Inches(0.4),
            font_size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

add_textbox(s, "MediaWale 2023 Pvt Ltd\nAll rights reserved · Proprietary & Confidential",
            Inches(8.6), Inches(4.65), Inches(4.2), Inches(0.65),
            font_size=9, color=MID_GREY, align=PP_ALIGN.CENTER)

# Quote
add_textbox(s,
    '"The MP who knows their constituency best — wins.\nNeedle makes that MP you."',
    Inches(0.6), Inches(6.5), Inches(7.5), Inches(0.7),
    font_size=13, italic=True, color=GOLD)

accent_bar(s, SLIDE_H - Inches(0.08), Inches(0.08))


# ─────────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────────
OUT = "/home/user/compass-needle/needle_pitch_deck.pptx"
prs.save(OUT)
print(f"Saved: {OUT}")
