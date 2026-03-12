"""
Needle — Political Buyer Pitch Deck
Target: MPs, MLAs, Councillors, Aspirants
needle_political_pitch.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─────────────────────────────────────
# PALETTE  — political energy meets tech
# ─────────────────────────────────────
NAVY        = RGBColor(0x07, 0x10, 0x1E)   # near-black navy
DEEP_NAVY   = RGBColor(0x0B, 0x17, 0x2B)
CARD        = RGBColor(0x12, 0x22, 0x35)
CARD2       = RGBColor(0x18, 0x2D, 0x42)
SAFFRON     = RGBColor(0xFF, 0x7A, 0x00)   # India saffron
GOLD        = RGBColor(0xF5, 0xC5, 0x18)   # election gold
INDIA_GREEN = RGBColor(0x13, 0x8A, 0x3C)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT       = RGBColor(0xD4, 0xE1, 0xED)
MID         = RGBColor(0x7A, 0x98, 0xB2)
RED         = RGBColor(0xE0, 0x30, 0x30)
TEAL        = RGBColor(0x0E, 0xB5, 0x9A)
PURPLE      = RGBColor(0x7C, 0x3A, 0xED)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]


# ─────────────────────────────────────
# LOW-LEVEL HELPERS
# ─────────────────────────────────────

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, color, line=False):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    if line:
        s.line.color.rgb = color
    else:
        s.line.fill.background()
    return s

def tx(slide, text, l, t, w, h,
       size=14, bold=False, italic=False,
       color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    r   = p.add_run()
    r.text = text
    r.font.name  = "Calibri"
    r.font.size  = Pt(size)
    r.font.bold  = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return box

def multiline(slide, lines, l, t, w, h, default_size=12, default_color=WHITE):
    """lines = list of (text, size, bold, color)"""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    for i, (text, size, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = text
        r.font.name  = "Calibri"
        r.font.size  = Pt(size or default_size)
        r.font.bold  = bold
        r.font.color.rgb = color or default_color
    return box

def top_bar(slide, color=SAFFRON, h=Inches(0.07)):
    rect(slide, 0, 0, W, h, color)

def bot_bar(slide, color=SAFFRON, h=Inches(0.07)):
    rect(slide, 0, H - h, W, h, color)

def label(slide, text, l, t, w=Inches(4), h=Inches(0.38),
          bg_c=SAFFRON, fg_c=NAVY):
    rect(slide, l, t, w, h, bg_c)
    tx(slide, text, l, t + Inches(0.05), w, h - Inches(0.05),
       size=9, bold=True, color=fg_c, align=PP_ALIGN.CENTER)

def divider(slide, t, color=CARD2, h=Inches(0.015)):
    rect(slide, Inches(0.6), t, W - Inches(1.2), h, color)

def card(slide, l, t, w, h, color=CARD):
    return rect(slide, l, t, w, h, color)

def brand(slide):
    tx(slide, "NEEDLE  ·  needle.mediawale.in  ·  Proprietary & Confidential",
       0, H - Inches(0.48), W - Inches(0.3), Inches(0.35),
       size=8, color=MID, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — HERO
# "The politician who knows their constituency best — wins."
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, NAVY)

# Full-bleed saffron stripe top
rect(s, 0, 0, W, Inches(0.5), SAFFRON)
# India green bottom stripe
rect(s, 0, H - Inches(0.5), W, Inches(0.5), INDIA_GREEN)

# Background texture — diagonal stripes effect
for i in range(8):
    x = Inches(7.0 + i * 0.9)
    rect(s, x, 0, Inches(0.04), H, RGBColor(0x0F, 0x1C, 0x30))

# Giant watermark text
tx(s, "NEEDLE", Inches(0.3), Inches(0.8), W, Inches(5),
   size=220, bold=True, color=RGBColor(0x0C, 0x1A, 0x2E))

# Label
label(s, "FOR MPs · MLAs · COUNCILLORS · ASPIRANTS", Inches(0.6), Inches(1.35),
      Inches(6.5), Inches(0.38), bg_c=SAFFRON, fg_c=NAVY)

# Main headline
tx(s, "The politician who knows\ntheir constituency best",
   Inches(0.6), Inches(1.9), Inches(9), Inches(2.1),
   size=52, bold=True, color=WHITE)
tx(s, "— wins.", Inches(0.6), Inches(3.9), Inches(9), Inches(0.9),
   size=52, bold=True, color=GOLD)

# Sub
tx(s, "Needle is your AI-powered political command centre.\nEvery grievance tracked. Every vote protected.",
   Inches(0.6), Inches(4.95), Inches(8), Inches(0.95),
   size=16, color=LIGHT)

brand(s)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — THE MATH OF RE-ELECTION
# Hard numbers on what determines winners
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, NAVY)
top_bar(s)
bot_bar(s, INDIA_GREEN)

label(s, "WHY THIS MATTERS NOW", Inches(0.6), Inches(0.22), Inches(3.2), Inches(0.36),
      bg_c=SAFFRON)

tx(s, "Elections are won or lost in\nthe years between elections.",
   Inches(0.6), Inches(0.75), Inches(8.5), Inches(1.5),
   size=38, bold=True, color=WHITE)

divider(s, Inches(2.35))

# Big stat cards
stats = [
    ("67%", "voter turnout in 2024\ngeneral election", SAFFRON),
    ("1.5L+", "average margin in a\ncompetitive Lok Sabha seat", GOLD),
    ("2.4M", "people every MP\nrepresents on average", TEAL),
    ("70%", "of voters say local\nissues decide their vote", INDIA_GREEN),
]

sw = Inches(2.8)
sh = Inches(2.4)
sx0 = Inches(0.55)
sy = Inches(2.55)

for i, (val, lbl, accent) in enumerate(stats):
    sx = sx0 + i * (sw + Inches(0.3))
    card(s, sx, sy, sw, sh)
    rect(s, sx, sy, sw, Inches(0.06), accent)
    tx(s, val, sx, sy + Inches(0.18), sw, Inches(0.9),
       size=42, bold=True, color=accent, align=PP_ALIGN.CENTER)
    tx(s, lbl, sx, sy + Inches(1.1), sw, Inches(1.1),
       size=11, color=LIGHT, align=PP_ALIGN.CENTER)

# The insight
rect(s, Inches(0.5), Inches(5.15), W - Inches(1.0), Inches(1.5), CARD2)
rect(s, Inches(0.5), Inches(5.15), Inches(0.07), Inches(1.5), GOLD)
tx(s, "The insight:", Inches(0.72), Inches(5.28), Inches(1.5), Inches(0.4),
   size=10, bold=True, color=GOLD)
tx(s, "Constituency responsiveness is the single most predictable indicator of re-election margin. "
      "MPs who track and resolve grievances publicly outperform those who don't — by 12–18 percentage "
      "points on satisfaction surveys. Satisfied voters turn out. Unsatisfied voters switch.",
   Inches(0.72), Inches(5.65), W - Inches(1.4), Inches(0.85),
   size=11, color=LIGHT)

brand(s)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — THE BRUTAL TRUTH
# "You represent 2.4 million people. With a team of 5."
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, NAVY)
top_bar(s, SAFFRON)
bot_bar(s, INDIA_GREEN)

label(s, "THE REALITY CHECK", Inches(0.6), Inches(0.22), Inches(2.8), Inches(0.36),
      bg_c=RED, fg_c=WHITE)

tx(s, "You represent 2.4 million people.\nWith a team of 5.",
   Inches(0.6), Inches(0.75), Inches(10), Inches(1.6),
   size=42, bold=True, color=WHITE)

divider(s, Inches(2.45))

# Left column — what's arriving daily
lw = Inches(5.6)
card(s, Inches(0.5), Inches(2.6), lw, Inches(4.2))
rect(s, Inches(0.5), Inches(2.6), lw, Inches(0.5), RED)
tx(s, "What arrives at your office every day",
   Inches(0.6), Inches(2.68), lw - Inches(0.2), Inches(0.36),
   size=11, bold=True, color=WHITE)

daily = [
    ("30–100",    "WhatsApp messages from constituents"),
    ("20–50",     "Letters and applications"),
    ("10–30",     "Phone calls your PA can't pick up"),
    ("5–15",      "Walk-ins waiting at your office"),
    ("3–10",      "Referrals forwarded by party workers"),
    ("Countless", "Social media tags and comments"),
]
dy = Inches(3.25)
for vol, desc in daily:
    tx(s, vol, Inches(0.65), dy, Inches(1.0), Inches(0.4),
       size=13, bold=True, color=RED)
    tx(s, desc, Inches(1.7), dy, Inches(4.2), Inches(0.38),
       size=11, color=LIGHT)
    dy += Inches(0.52)

# Right column — what happens to most of it
rx = Inches(6.6)
card(s, rx, Inches(2.6), lw, Inches(4.2))
rect(s, rx, Inches(2.6), lw, Inches(0.5), RGBColor(0x80, 0x30, 0x00))
tx(s, "What actually happens to most of it",
   rx + Inches(0.1), Inches(2.68), lw - Inches(0.2), Inches(0.36),
   size=11, bold=True, color=WHITE)

fates = [
    ("Forgotten", "in a WhatsApp thread no one revisits"),
    ("Mis-routed", "to the wrong department or constituency"),
    ("Duplicated", "same complaint filed 3x by the same family"),
    ("Untracked", "no follow-up, no resolution date, no record"),
    ("Politically costly", "constituent tells 10 neighbours you ignored them"),
    ("Lost forever", "when your PA changes"),
]
fy = Inches(3.25)
for verdict, desc in fates:
    tx(s, verdict, rx + Inches(0.15), fy, Inches(1.6), Inches(0.4),
       size=12, bold=True, color=RGBColor(0xFF, 0x88, 0x55))
    tx(s, desc, rx + Inches(1.85), fy, Inches(4.5), Inches(0.38),
       size=11, color=LIGHT)
    fy += Inches(0.52)

brand(s)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — WHAT YOUR CONSTITUENT ACTUALLY WANTS
# Empathy slide — their journey
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, NAVY)
top_bar(s, SAFFRON)
bot_bar(s, INDIA_GREEN)

label(s, "THE CONSTITUENT JOURNEY", Inches(0.6), Inches(0.22), Inches(3.2), Inches(0.36))

tx(s, "Ramesh in Nagpur has one problem.\nHe just wants it fixed.",
   Inches(0.6), Inches(0.75), Inches(10), Inches(1.4),
   size=38, bold=True, color=WHITE)

divider(s, Inches(2.2))

# Journey stages with emotion
stages = [
    ("The Problem",        SAFFRON,  "Ramesh's street has had no streetlight for 8 months. His daughter walks home from tuition in the dark. He has complained to the municipality 3 times. Nothing happened."),
    ("He Finds You",       GOLD,     "A neighbour says 'Message the MP on WhatsApp.' Ramesh finds the number. He types in Marathi. He has never done this before. He waits."),
    ("The Make-or-Break",  TEAL,     "THIS is the moment. The next 3 minutes decide whether Ramesh becomes your advocate or your opponent. Does he get an answer — or silence?"),
    ("With Needle",        INDIA_GREEN, "Needle reads his message in Marathi. AI acknowledges instantly. Case is logged, tagged Infrastructure > Lighting, location pinned to his ward. You see it on your dashboard. You act."),
    ("6 Weeks Later",      RGBColor(0xFF,0xFF,0xFF), "The streetlight is fixed. Ramesh got updates twice. He tells 40 people: 'Message the MP — they actually respond.' That's 40 potential votes protected."),
]

sw2 = Inches(2.3)
sh2 = Inches(3.8)
sx2 = Inches(0.45)
sy2 = Inches(2.35)

for i, (stage, accent, story) in enumerate(stages):
    sxc = sx2 + i * (sw2 + Inches(0.2))
    card(s, sxc, sy2, sw2, sh2)
    rect(s, sxc, sy2, sw2, Inches(0.06), accent)
    tx(s, f"{i+1}", sxc, sy2 + Inches(0.12), sw2, Inches(0.45),
       size=18, bold=True, color=accent, align=PP_ALIGN.CENTER)
    tx(s, stage, sxc, sy2 + Inches(0.55), sw2, Inches(0.45),
       size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tx(s, story, sxc + Inches(0.1), sy2 + Inches(1.1),
       sw2 - Inches(0.2), Inches(2.5),
       size=9, color=LIGHT)

brand(s)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — INTRODUCING NEEDLE
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, NAVY)
top_bar(s, SAFFRON)
bot_bar(s, INDIA_GREEN)

label(s, "INTRODUCING NEEDLE", Inches(0.6), Inches(0.22), Inches(2.8), Inches(0.36),
      bg_c=GOLD, fg_c=NAVY)

tx(s, "Your 24/7 political command centre.\nBuilt for Indian democracy.",
   Inches(0.6), Inches(0.75), Inches(9), Inches(1.4),
   size=38, bold=True, color=WHITE)

divider(s, Inches(2.2))

# Central feature grid
features = [
    ("WhatsApp\nInbox", SAFFRON,
     "Citizens message you in any Indian language. Needle reads, categorises, and routes every single one — while you sleep."),
    ("AI Drafter",  GOLD,
     "Draft letters, Parliament questions, Private Member Bills, and speeches in the official Government of India format. In seconds."),
    ("Case\nTracker", TEAL,
     "Every complaint gets a case file — status, category, location, history. Nothing falls through the cracks. Ever."),
    ("CSR\nFinder",   INDIA_GREEN,
     "Discover which corporates have CSR funds available for your constituency. Generate ready-to-send proposals automatically."),
    ("Document\nCopilot", PURPLE,
     "Upload any government report, scheme notification, or budget. Ask it questions in Hindi or English. Get cited answers instantly."),
    ("News\nIntel",  RGBColor(0xF0, 0x60, 0x60),
     "A live news feed filtered to your constituency and state. Know what's happening in your area before your opponent does."),
]

fw = Inches(3.8)
fh = Inches(2.1)
fx0 = Inches(0.5)
fy0 = Inches(2.4)

for i, (name, accent, desc) in enumerate(features):
    col = i % 3
    row = i // 3
    fx = fx0 + col * (fw + Inches(0.3))
    fy = fy0 + row * (fh + Inches(0.2))
    card(s, fx, fy, fw, fh)
    rect(s, fx, fy, Inches(0.06), fh, accent)
    tx(s, name, fx + Inches(0.2), fy + Inches(0.15),
       fw - Inches(0.3), Inches(0.6),
       size=14, bold=True, color=WHITE)
    tx(s, desc, fx + Inches(0.2), fy + Inches(0.75),
       fw - Inches(0.3), Inches(1.2),
       size=10, color=LIGHT)

brand(s)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — YOUR PA SLEEPS. NEEDLE DOESN'T.
# WhatsApp engine deep-dive
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, NAVY)
top_bar(s, SAFFRON)
bot_bar(s, INDIA_GREEN)

label(s, "THE WHATSAPP ENGINE", Inches(0.6), Inches(0.22), Inches(2.8), Inches(0.36),
      bg_c=SAFFRON)

tx(s, "Your PA sleeps. Needle doesn't.",
   Inches(0.6), Inches(0.75), Inches(9), Inches(0.9),
   size=42, bold=True, color=WHITE)
tx(s, "Every WhatsApp message from every constituent — received, read, categorised, and acknowledged.\nAutomatically. In their language. At 3am if needed.",
   Inches(0.6), Inches(1.65), Inches(9.5), Inches(0.85),
   size=14, color=LIGHT)

divider(s, Inches(2.6))

# Left: language/capability callouts
lang_cards = [
    ("Hindi",       "हिंदी", SAFFRON),
    ("Marathi",     "मराठी", GOLD),
    ("Kannada",     "ಕನ್ನಡ", TEAL),
    ("Hinglish",    "Mix", INDIA_GREEN),
    ("English",     "EN", PURPLE),
]
lx = Inches(0.5)
ly = Inches(2.8)
for i, (lang, script, c) in enumerate(lang_cards):
    lxi = lx + i * Inches(1.1)
    rect(s, lxi, ly, Inches(1.0), Inches(0.85), CARD2)
    rect(s, lxi, ly, Inches(1.0), Inches(0.04), c)
    tx(s, script, lxi, ly + Inches(0.05), Inches(1.0), Inches(0.4),
       size=16, bold=True, color=c, align=PP_ALIGN.CENTER)
    tx(s, lang, lxi, ly + Inches(0.48), Inches(1.0), Inches(0.3),
       size=8, color=LIGHT, align=PP_ALIGN.CENTER)

# Flow
flow_steps = [
    ("Citizen sends\nWhatsApp", "In any language,\nany time of day"),
    ("AI reads\n& classifies", "Category + Location\n+ Urgency level"),
    ("Auto-reply\nsent", "In citizen's language\nwithin 3 seconds"),
    ("Case logged\non dashboard", "Status, tags, history\nfully tracked"),
    ("MP takes\naction", "Assign, escalate,\nor resolve"),
]

fw2 = Inches(2.2)
fh2 = Inches(1.8)
fx2 = Inches(0.45)
fy2 = Inches(3.85)

for i, (step, sub) in enumerate(flow_steps):
    fxi = fx2 + i * (fw2 + Inches(0.35))
    is_ai = (i == 1)
    bg_c = SAFFRON if is_ai else CARD
    tc = NAVY if is_ai else WHITE
    sc = NAVY if is_ai else MID
    card(s, fxi, fy2, fw2, fh2, bg_c)
    tx(s, step, fxi, fy2 + Inches(0.2), fw2, Inches(0.7),
       size=13, bold=True, color=tc, align=PP_ALIGN.CENTER)
    tx(s, sub, fxi, fy2 + Inches(0.95), fw2, Inches(0.7),
       size=9, color=sc, align=PP_ALIGN.CENTER)
    if i < len(flow_steps) - 1:
        tx(s, "▶", fxi + fw2 + Inches(0.08), fy2 + Inches(0.65),
           Inches(0.22), Inches(0.4), size=14, color=SAFFRON, align=PP_ALIGN.CENTER)

brand(s)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — OUTPERFORM IN PARLIAMENT
# AI Drafter for MPs
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, NAVY)
top_bar(s, SAFFRON)
bot_bar(s, INDIA_GREEN)

label(s, "PARLIAMENT PERFORMANCE", Inches(0.6), Inches(0.22), Inches(3.0), Inches(0.36),
      bg_c=GOLD, fg_c=NAVY)

tx(s, "Your party notices who asks\nthe sharpest questions.",
   Inches(0.6), Inches(0.75), Inches(8.5), Inches(1.4),
   size=38, bold=True, color=WHITE)

divider(s, Inches(2.2))

# Document types grid — 2x2
docs = [
    ("Official Letters",         SAFFRON,
     "GoI standard format. F.No., formal Rajbhasha Hindi. Write to any Minister or Department in under 2 minutes. Assertive, diplomatic, or formal tone — your choice."),
    ("Parliament Questions",      GOLD,
     "Starred and Unstarred questions for Lok Sabha and Rajya Sabha. Proper format, sharp subject lines, data-backed. Ask more. Ask better."),
    ("Private Member Bills",      TEAL,
     "Clause-by-clause bill drafts with Statement of Objects and Reasons. Turn constituency problems into legislative proposals."),
    ("Speeches & Statements",     INDIA_GREEN,
     "Budget reply speeches, committee statements, constituency development addresses. With data, with arguments, with impact."),
]

dw = Inches(5.8)
dh = Inches(2.1)
dx0 = Inches(0.5)
dy0 = Inches(2.45)

for i, (dtype, accent, desc) in enumerate(docs):
    col = i % 2
    row = i // 2
    dx = dx0 + col * (dw + Inches(0.55))
    dy = dy0 + row * (dh + Inches(0.2))
    card(s, dx, dy, dw, dh)
    rect(s, dx, dy, dw, Inches(0.06), accent)
    tx(s, dtype, dx + Inches(0.18), dy + Inches(0.18),
       dw - Inches(0.3), Inches(0.45),
       size=15, bold=True, color=WHITE)
    tx(s, desc, dx + Inches(0.18), dy + Inches(0.72),
       dw - Inches(0.3), Inches(1.2),
       size=10, color=LIGHT)

# DNA style callout
rect(s, Inches(0.5), Inches(6.8), W - Inches(1.0), Inches(0.42), CARD2)
tx(s, "DNA Style Templates — Needle learns your writing voice. Every draft sounds like you, not a robot.",
   Inches(0.65), Inches(6.88), W - Inches(1.3), Inches(0.32),
   size=10, color=GOLD)

brand(s)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — BRING CRORES HOME
# CSR & Fund Intelligence
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, NAVY)
top_bar(s, SAFFRON)
bot_bar(s, INDIA_GREEN)

label(s, "CSR & FUNDS INTELLIGENCE", Inches(0.6), Inches(0.22), Inches(3.2), Inches(0.36),
      bg_c=INDIA_GREEN, fg_c=WHITE)

tx(s, "Bring crores to your constituency.\nWithout spending a rupee of government funds.",
   Inches(0.6), Inches(0.75), Inches(10), Inches(1.4),
   size=36, bold=True, color=WHITE)

divider(s, Inches(2.2))

# Big number left
rect(s, Inches(0.5), Inches(2.45), Inches(4.2), Inches(4.2), CARD)
rect(s, Inches(0.5), Inches(2.45), Inches(4.2), Inches(0.07), INDIA_GREEN)

tx(s, "₹25,000 Cr+",
   Inches(0.5), Inches(2.7), Inches(4.2), Inches(1.0),
   size=38, bold=True, color=INDIA_GREEN, align=PP_ALIGN.CENTER)
tx(s, "unspent corporate CSR funds\nin India every year",
   Inches(0.5), Inches(3.65), Inches(4.2), Inches(0.65),
   size=12, color=LIGHT, align=PP_ALIGN.CENTER)

divider(s, Inches(4.4), color=CARD2)

tx(s, "Most MPs never access this money because they don't know which companies are eligible, what they fund, or how to apply.",
   Inches(0.65), Inches(4.55), Inches(3.9), Inches(1.1),
   size=10, color=LIGHT)

# Right column — what needle does
rx = Inches(5.1)
rw = Inches(7.7)

csr_steps = [
    (SAFFRON, "Gap Analysis",
     "Needle reads your grievance data — your top complaint categories ARE your constituency's unmet needs. Water? Roads? Schools? It knows."),
    (GOLD, "Company Matching",
     "Cross-references complaint categories with companies that have CSR budgets in matching sectors. Finds the fit automatically."),
    (INDIA_GREEN, "Proposal Generator",
     "One click generates a formatted CSR proposal — project scope, budget, impact metrics, SDG alignment, Schedule VII compliance. Ready to send."),
    (TEAL, "Projects Menu",
     "10+ pre-costed constituency projects: solar streetlights (₹25K), borewells (₹1.5L), RO plants (₹4.5L), smart classrooms (₹2.5L), health sub-centres (₹15L)."),
]

py2 = Inches(2.45)
for accent, title, desc in csr_steps:
    card(s, rx, py2, rw, Inches(0.95))
    rect(s, rx, py2, Inches(0.06), Inches(0.95), accent)
    tx(s, title, rx + Inches(0.18), py2 + Inches(0.08),
       Inches(2.0), Inches(0.38),
       size=12, bold=True, color=WHITE)
    tx(s, desc, rx + Inches(2.3), py2 + Inches(0.08),
       rw - Inches(2.5), Inches(0.75),
       size=10, color=LIGHT)
    py2 += Inches(1.05)

brand(s)


# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — BUILT FOR EVERY LEVEL
# MP / MLA / Councillor / Aspirant — one slide, four segments
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, NAVY)
top_bar(s, SAFFRON)
bot_bar(s, INDIA_GREEN)

label(s, "MADE FOR YOUR ROLE", Inches(0.6), Inches(0.22), Inches(2.6), Inches(0.36),
      bg_c=SAFFRON)

tx(s, "Whether you hold a seat or want one —\nNeedle is built for you.",
   Inches(0.6), Inches(0.75), Inches(10), Inches(1.3),
   size=36, bold=True, color=WHITE)

divider(s, Inches(2.1))

segments = [
    ("MP\n(Lok Sabha / Rajya Sabha)", SAFFRON, NAVY,
     [
         "Parliament Questions: Starred, Unstarred, Short Notice",
         "Private Member Bill drafting with full clause format",
         "Ministry fund tracking — know where your constituency's share goes",
         "Committee work support and brief preparation",
         "Sansad TV case tracking (SansadX integration)",
         "Constituency grievance dashboard across all Assembly segments",
     ]),
    ("MLA\n(State Assembly)", GOLD, NAVY,
     [
         "State-level scheme delivery tracking (MGNREGA, PMAY, Jal Jeevan)",
         "Vidhansabha question drafting in state language + Hindi",
         "MLAS fund project pipeline and CSR matching",
         "Constituency development: road, water, health complaints",
         "Booth-level grievance mapping with geography engine",
         "State budget analysis and response speech drafting",
     ]),
    ("Municipal Councillor /\nPanchayat Member", TEAL, WHITE,
     [
         "Ward-level complaint tracker: roads, drains, lights, water supply",
         "Grievance-to-action pipeline with WhatsApp intake",
         "Simple case dashboard: pending → in-progress → resolved",
         "Direct document drafting: letters to BMC / corporation / collector",
         "CSR micro-projects for ward (₹25K – ₹5L scale)",
         "Ward health score: % issues resolved per month",
     ]),
    ("Election Aspirant /\nCandidate-in-Waiting", PURPLE, WHITE,
     [
         "Start listening to your constituency NOW — before you win",
         "Build a database of citizen problems by booth and ward",
         "Understand top 5 issues before your first campaign speech",
         "Draft your manifesto from real complaint data — not assumptions",
         "Establish a WhatsApp presence as 'the candidate who cares'",
         "Know the CSR landscape: promise real projects with real funding",
     ]),
]

sw3 = Inches(2.9)
sh3 = Inches(4.2)
sx3 = Inches(0.45)
sy3 = Inches(2.25)

for i, (role, accent, text_c, bullets) in enumerate(segments):
    sxc = sx3 + i * (sw3 + Inches(0.25))
    card(s, sxc, sy3, sw3, sh3)
    rect(s, sxc, sy3, sw3, Inches(0.55), accent)
    tx(s, role, sxc, sy3 + Inches(0.05), sw3, Inches(0.5),
       size=11, bold=True, color=text_c, align=PP_ALIGN.CENTER)
    by3 = sy3 + Inches(0.72)
    for bullet in bullets:
        tx(s, "›  " + bullet, sxc + Inches(0.1), by3,
           sw3 - Inches(0.15), Inches(0.45),
           size=9, color=LIGHT)
        by3 += Inches(0.55)

brand(s)


# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — YOUR OPPONENT ISN'T SLEEPING
# Competitive / urgency slide
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, NAVY)
top_bar(s, RED)
bot_bar(s, INDIA_GREEN)

label(s, "THE COMPETITIVE REALITY", Inches(0.6), Inches(0.22), Inches(3.0), Inches(0.36),
      bg_c=RED, fg_c=WHITE)

tx(s, "Your opponent isn't sleeping.\nAre you ahead — or behind?",
   Inches(0.6), Inches(0.75), Inches(9), Inches(1.4),
   size=38, bold=True, color=WHITE)

divider(s, Inches(2.2))

# Comparison table
table_data = [
    ("CAPABILITY",                          "Without Needle",  "With Needle"),
    ("Grievance intake after office hours",  "❌  Nothing",      "✅  24/7 WhatsApp AI"),
    ("Complaint categorised automatically",  "❌  PA does it",    "✅  GPT-4 in < 3 seconds"),
    ("Multi-language constituent support",   "❌  Hindi only",    "✅  5 languages + dialects"),
    ("Parliament question drafted in 1 min", "❌  2–4 hours",     "✅  60 seconds"),
    ("CSR fund matching for constituency",   "❌  Personal contacts","✅  Automated matching"),
    ("Constituency news filtered to your area","❌  Generic news", "✅  Constituency-aware feed"),
    ("Case history per constituent",         "❌  Scattered notes","✅  Full case database"),
    ("Document analysis (PDF Q&A)",          "❌  Read it yourself","✅  AI Copilot"),
    ("Re-election advantage",                "❌  Hope",          "✅  Data + Responsiveness"),
]

col_widths = [Inches(4.5), Inches(3.3), Inches(4.0)]
tx0 = [Inches(0.5), Inches(5.1), Inches(8.5)]
ty0 = Inches(2.42)
row_h = Inches(0.43)

for ri, row in enumerate(table_data):
    row_y = ty0 + ri * row_h
    is_header = (ri == 0)
    row_bg = CARD2 if ri % 2 == 0 and not is_header else CARD if not is_header else RGBColor(0x10, 0x25, 0x3A)
    rect(s, Inches(0.5), row_y, sum(col_widths), row_h, row_bg)

    for ci, (cell_text, cw, cx) in enumerate(zip(row, col_widths, tx0)):
        c_color = GOLD if is_header else (RED if "❌" in cell_text else INDIA_GREEN if "✅" in cell_text else LIGHT)
        sz = 9 if not is_header else 10
        tx(s, cell_text, cx + Inches(0.1), row_y + Inches(0.07), cw - Inches(0.15), row_h - Inches(0.08),
           size=sz, bold=is_header, color=c_color)

brand(s)


# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — THE NUMBERS THAT MATTER TO YOU
# Political ROI — vote math
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, NAVY)
top_bar(s, SAFFRON)
bot_bar(s, INDIA_GREEN)

label(s, "YOUR POLITICAL ROI", Inches(0.6), Inches(0.22), Inches(2.6), Inches(0.36),
      bg_c=GOLD, fg_c=NAVY)

tx(s, "This is not a cost.\nIt's the cheapest vote you'll ever protect.",
   Inches(0.6), Inches(0.75), Inches(9), Inches(1.4),
   size=36, bold=True, color=WHITE)

divider(s, Inches(2.2))

# Left: the vote math
rect(s, Inches(0.5), Inches(2.45), Inches(5.7), Inches(4.3), CARD)
rect(s, Inches(0.5), Inches(2.45), Inches(5.7), Inches(0.06), GOLD)

tx(s, "The Vote Maths", Inches(0.65), Inches(2.55), Inches(5.4), Inches(0.4),
   size=13, bold=True, color=GOLD)

vote_lines = [
    ("1 resolved complaint  =  1 grateful family", LIGHT, 11, False),
    ("1 family  =  avg. 4.2 registered voters", LIGHT, 11, False),
    ("1 family shares story with 3 neighbours", LIGHT, 11, False),
    ("", LIGHT, 6, False),
    ("50 resolved cases/month", WHITE, 13, True),
    ("→  50 families  →  210 direct votes protected", GOLD, 13, True),
    ("→  150 neighbour households influenced", GOLD, 11, False),
    ("→  ~630 additional votes in your orbit", GOLD, 11, False),
    ("", LIGHT, 6, False),
    ("Needle costs ₹10,000/month.", LIGHT, 11, False),
    ("That's ₹12 per protected vote.", WHITE, 13, True),
    ("Your last election campaign spent ₹120–400+\nper voter reached.", MID, 10, False),
]

ly2 = Inches(3.05)
for text, color, size, bold in vote_lines:
    if text:
        tx(s, text, Inches(0.65), ly2, Inches(5.3), Inches(0.42),
           size=size, bold=bold, color=color)
    ly2 += Inches(0.36) if size >= 11 else Inches(0.2)

# Right: key stat boxes
rx2 = Inches(6.7)
rw2 = Inches(6.1)

stat_boxes = [
    ("₹12", "per protected vote with Needle", GOLD),
    ("₹120–400+", "per voter reached in traditional campaign", RED),
    ("10x cheaper", "than booth-level campaign spend per contact", INDIA_GREEN),
    ("24/7", "availability — your PA works 8 hours", TEAL),
]

sy4 = Inches(2.45)
for val, lbl, accent in stat_boxes:
    card(s, rx2, sy4, rw2, Inches(0.9))
    rect(s, rx2, sy4, Inches(0.06), Inches(0.9), accent)
    tx(s, val, rx2 + Inches(0.2), sy4 + Inches(0.1), Inches(2.2), Inches(0.5),
       size=22, bold=True, color=accent)
    tx(s, lbl, rx2 + Inches(2.5), sy4 + Inches(0.18), rw2 - Inches(2.7), Inches(0.55),
       size=10, color=LIGHT)
    sy4 += Inches(1.05)

brand(s)


# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — FOR THE ASPIRANT
# Special slide for candidates who haven't won yet
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, NAVY)
top_bar(s, PURPLE)
bot_bar(s, INDIA_GREEN)

label(s, "FOR THE ASPIRANT", Inches(0.6), Inches(0.22), Inches(2.6), Inches(0.36),
      bg_c=PURPLE, fg_c=WHITE)

tx(s, "Win the election before\nthe voting booth opens.",
   Inches(0.6), Inches(0.75), Inches(9), Inches(1.4),
   size=38, bold=True, color=WHITE)

divider(s, Inches(2.2))

tx(s, "The candidates who win in 2026 and 2029 are building constituency intelligence TODAY.",
   Inches(0.6), Inches(2.3), Inches(12), Inches(0.5),
   size=14, italic=True, color=GOLD)

# Three aspirant use-cases
asp_cards = [
    ("Know Your\nConstituency",   PURPLE,
     "Start a WhatsApp number for your constituency RIGHT NOW. Collect real grievances from real people. Know the top 10 issues in every ward before you file your nomination. Your opponent is guessing. You're not.",
     ["Complaint data by booth and ward", "Language-wise constituent profile", "Top 10 issues per Assembly segment"]),
    ("Build Your\nNarrative",      GOLD,
     "Your election manifesto shouldn't be promises — it should be solutions. Needle turns real grievance data into a data-backed manifesto. 'My constituency told me their top problem is water. Here's my plan.' That's not a speech — that's a movement.",
     ["Manifesto from real data, not party templates", "CSR project commitments with actual company names", "Draft speeches per constituency visit"]),
    ("Establish\nCredibility",      TEAL,
     "Use Needle to respond to every WhatsApp message, show you care before you win, and document your responsiveness. When election season comes, you don't introduce yourself — you show a 12-month record of being the only candidate who actually replied.",
     ["12-month response record", "Documented constituent interactions", "Digital proof of accountability"]),
]

aw = Inches(3.8)
ah = Inches(4.3)
ax0 = Inches(0.5)
ay0 = Inches(2.9)

for i, (title, accent, story, bullets) in enumerate(asp_cards):
    ax = ax0 + i * (aw + Inches(0.35))
    card(s, ax, ay0, aw, ah)
    rect(s, ax, ay0, aw, Inches(0.06), accent)
    tx(s, title, ax, ay0 + Inches(0.15), aw, Inches(0.55),
       size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tx(s, story, ax + Inches(0.15), ay0 + Inches(0.85),
       aw - Inches(0.3), Inches(1.85),
       size=9, color=LIGHT)
    by4 = ay0 + Inches(2.8)
    for b in bullets:
        tx(s, "✓  " + b, ax + Inches(0.15), by4,
           aw - Inches(0.3), Inches(0.38),
           size=9, color=accent, bold=True)
        by4 += Inches(0.44)

brand(s)


# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — WHAT NEEDLE KNOWS ABOUT YOUR CONSTITUENCY
# Dramatic "intelligence" demo slide
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, NAVY)
top_bar(s, SAFFRON)
bot_bar(s, INDIA_GREEN)

label(s, "THE INTELLIGENCE ADVANTAGE", Inches(0.6), Inches(0.22), Inches(3.3), Inches(0.36),
      bg_c=TEAL, fg_c=NAVY)

tx(s, "While others campaign on\nassumptions, you govern on data.",
   Inches(0.6), Inches(0.75), Inches(9.5), Inches(1.4),
   size=36, bold=True, color=WHITE)

divider(s, Inches(2.2))

# Sample "live dashboard" mockup as text cards
tx(s, "Sample constituency dashboard — what Needle shows you every morning:",
   Inches(0.6), Inches(2.3), Inches(11), Inches(0.38),
   size=10, italic=True, color=MID)

# Dashboard cards
dash_data = [
    ("TOP GRIEVANCE THIS WEEK",   SAFFRON, "Water Supply",  "34 complaints · ↑ 18% vs last week · Highest in Bhosari ward"),
    ("PENDING CASES",              RED,     "127 open",      "18 critical · 43 in progress · 66 new this week"),
    ("RESOLVED THIS MONTH",        INDIA_GREEN, "89 cases",  "↑ 32% resolution rate · Avg. 4.2 days to close"),
    ("CSR OPPORTUNITY IDENTIFIED", GOLD,    "₹45L available","3 companies matched to your water + road gaps"),
    ("PARLIAMENT QUESTION DUE",    TEAL,    "2 questions",   "Drafter ready · Rail Budget + MGNREGA allocation"),
    ("CONSTITUENCY NEWS ALERTS",   PURPLE,  "4 new alerts",  "Local dam dispute · New highway announcement · School strike"),
]

dw2 = Inches(5.8)
dh2 = Inches(0.78)
dx0b = Inches(0.5)
dy0b = Inches(2.82)

for i, (label_t, accent, main, sub) in enumerate(dash_data):
    col = i % 2
    row = i // 2
    dx2 = dx0b + col * (dw2 + Inches(0.55))
    dy2 = dy0b + row * (dh2 + Inches(0.12))
    card(s, dx2, dy2, dw2, dh2, CARD)
    rect(s, dx2, dy2, Inches(0.06), dh2, accent)
    tx(s, label_t, dx2 + Inches(0.15), dy2 + Inches(0.05), Inches(2.4), Inches(0.28),
       size=7, bold=True, color=accent)
    tx(s, main, dx2 + Inches(0.15), dy2 + Inches(0.3), Inches(1.8), Inches(0.35),
       size=16, bold=True, color=WHITE)
    tx(s, sub, dx2 + Inches(2.0), dy2 + Inches(0.2), dw2 - Inches(2.2), Inches(0.48),
       size=9, color=LIGHT)

brand(s)


# ═══════════════════════════════════════════════════════════════
# SLIDE 14 — THE CALL TO ACTION
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, NAVY)

# Bold saffron top stripe
rect(s, 0, 0, W, Inches(0.55), SAFFRON)
# India green bottom stripe
rect(s, 0, H - Inches(0.55), W, Inches(0.55), INDIA_GREEN)

# Background large text
tx(s, "START", Inches(-0.2), Inches(0.4), W, Inches(4.5),
   size=250, bold=True, color=RGBColor(0x0C, 0x1A, 0x2E))

# Headline
tx(s, "Your constituency is waiting.\nStart today.",
   Inches(0.6), Inches(0.8), Inches(9), Inches(1.7),
   size=46, bold=True, color=WHITE)

tx(s, "Live in 48 hours. No IT team needed. Cancel anytime.",
   Inches(0.6), Inches(2.5), Inches(8), Inches(0.55),
   size=16, color=LIGHT)

divider(s, Inches(3.15))

# Three action tracks
tracks = [
    ("For MPs & MLAs",        SAFFRON, NAVY,
     "Free 30-day pilot\nFull product, your constituency, zero commitment."),
    ("For Councillors",        GOLD, NAVY,
     "Starter plan from ₹5,000/mo\nWard-level WhatsApp intake + case tracker."),
    ("For Aspirants",          PURPLE, WHITE,
     "Aspirant Plan — coming Q2 2025\nJoin the waitlist. Build your base now."),
]

tw3 = Inches(3.8)
th3 = Inches(1.8)
tx3 = Inches(0.5)
ty3 = Inches(3.35)

for i, (who, accent, text_c, offer) in enumerate(tracks):
    txi = tx3 + i * (tw3 + Inches(0.35))
    card(s, txi, ty3, tw3, th3, accent)
    tx(s, who, txi, ty3 + Inches(0.12), tw3, Inches(0.45),
       size=14, bold=True, color=text_c, align=PP_ALIGN.CENTER)
    tx(s, offer, txi, ty3 + Inches(0.62), tw3, Inches(1.0),
       size=11, color=text_c, align=PP_ALIGN.CENTER)

# Contact
rect(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.0), CARD2)
tx(s, "WhatsApp:  +91 XXXXX XXXXX   ·   Email:  contact@mediawale.in   ·   needle.mediawale.in",
   Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.5),
   size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
tx(s, "Demo available — 30 minutes, your team, your constituency, your questions.",
   Inches(0.5), Inches(6.08), Inches(12.3), Inches(0.38),
   size=11, color=LIGHT, align=PP_ALIGN.CENTER)

# Bottom quote
tx(s, '"The booth agent who never asks for a party ticket."',
   Inches(0.6), Inches(6.6), Inches(12.0), Inches(0.5),
   size=13, italic=True, bold=False, color=SAFFRON, align=PP_ALIGN.CENTER)

brand(s)


# ─────────────────────────────────────
# SAVE
# ─────────────────────────────────────
OUT = "/home/user/compass-needle/needle_political_pitch.pptx"
prs.save(OUT)
print(f"Saved → {OUT}")
